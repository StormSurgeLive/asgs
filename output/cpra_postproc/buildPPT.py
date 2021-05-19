#!/usr/bin/env python
#####################################################
# AUTHORS: Matthew V Bilskie, PhD
#          Louisiana State University
#
#          Jason Fleming, PhD
#          Seahorse Coastal Consulting
# COPYRIGHT 2018
#
#####################################################
#
# TO RUN IN A PYTHON SHELL USING CMD ARGS
#import sys
#sys.argv = ['./buildPPT.py','FigureGenFilename']
#execfile('./buildPPT.py')
#
#####################################################

import sys
from datetime import datetime
import pytz
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

tz = pytz.timezone('America/Chicago')
today = datetime.now(tz)

# Get command line argument
fname = sys.argv[1]
rprop = sys.argv[2]

# File names and stations
fnames = ['WSE_IHNC02_USACE_76030.png',
         'WSE_IHNC01_USGS_073802332.png',
         'WSE_17StCanal_USACE_85625.png',
         'WSE_LakefrontAirport_USACE_85670.png',
         'WSE_Rigolets_USACE_85700.png',
         'WSE_Mandeville_USACE_85575.png',
         'WSE_LPV144_USACE_76010.png',
         'WSE_LPV149_USACE_85760.png',
         'WSE_BayouBienv_USACE_76025.png',
         'WSE_WBV90_USACE_76265.png',
         'WSE_HarveyCanalBoom_USACE_76230.png',
         'WSE_Lafitte_USACE_82875.png',
         'WSE_WBV162_USACE_82742.png',
         'WSE_WBV7274_USACE_82715.png',
         'WSE_WBV09b_USACE_82762.png',
         'WSE_NOV13_USGS_07380260.png',
         'WSE_BaraPass_USGS_073802516.png',
         'WSE_HoumaNavCanal_USACE_76305.png',
         'WSE_BayouSale_USACE_76560.png',
         'WSE_BayouBoeuf_USGS_073814675.png',
         'WSE_MorganCity_USACE_03780.png',
         'WSE_CalcRiv_USGS_08017118.png',
         'WSE_Venice_USACE_01480.png',
         'WSE_NOV14_USACE_01440.png',
         'WSE_WestPoint_USACE_01400.png',
         'WSE_Alliance_USACE_01390.png',
         'WSE_Carrollton_USACE_01300.png',
         'WSE_BCSpillwayN_USACE_01275.png',
         'WSE_Reserve_USACE_01260.png']

# Station names correspond to the order of fnames
staName = ['IHNC Surge Barrier East - Flood Side, LA (IHNC-02, CPRA) (76030, USACE)',
            'Seabrook Complex - Flood Side, LA (IHNC-01, CPRA) (073802332, USGS)',
            'Outfall: London Ave, Orleans Ave, 17th St. (17StCanal, CPRA) (85625, USACE)',
            'Lake Pontchartrain at Lakefront Airport, LA (LakefrontAirport, CPRA)\n(85670, USACE)',
            'Rigolets near Lake Pontchartrain, LA (Rigolets, CPRA) (85700, USACE)',
            'Lake Pontchartrain at Mandeville, LA (Mandeville, CPRA) (85575, USACE)',
            'Bayou Dupre Sector Gate - East/Flood Side, LA (LPV-144, CPRA)\n(76010, USACE)',
            'Caernarvon Canal Sector Gate - South/Flood Side, LA (LPV-149, CPRA)\n(85760, USACE)',
            'Bayou Bienvenue Floodgate, LA (BayouBienv, CPRA) (76025, USACE)',
            'GIWW at West Closure Complex - Flood Side, LA (WBV-90, CPRA)\n(76265, USACE)',
            'Harvey Canal at Boomtown Casino, LA (HarveyCanalBoom, CPRA)\n(76230, USACE)',
            'Barataria Waterway at Lafitte, LA (Lafitte, CPRA) (82875, USACE)',
            'Bayou Segnette Closure - Flood Side, LA (WBV-16.2, CPRA) (82742, USACE)',
            'Bayou Verret / W. Tie-In Sector Gate Flood Side, LA (WBV-72/74, CPRA)\n(82715, USACE)',
            'Hero Canal Stop-Log Gate - Flood Side/West, LA (WBV-09b, CPRA)\n(82762, USACE)',
            'MS River at Empire Floodgate, LA (NOV-13, CPRA) (07380260, USGS)',
            'Baratria Pass at Grand Isle, LA (BaraPass, CPRA) (073802516, USGS)',
            'Houma Navigation Canal (76305, USACE)',
            'GIWW at Bayou Sale (76560, USACE)',
            'Bayou Boeuf at Railroad Bridge (073814675, USGS)',
            'Lower Atchafalaya River at Morgan City (03780, USACE)',
            'Calcasieu River at Cameron, LA (CalcRiv, CPRA) (8017118, USGS)',
            'MS River at Venice (01480, USACE); RM 10-11',
            'MS River at Empire Lock, LA (NOV-14, CPRA) (01440, USACE);\nRM 29-30',
            'MS River at West Point a la Hache (01400, USACE); RM 48-49',
            'MS River at Alliance (01390, USACE); RM 62-63',
            'MS River at Carrollton (01300, USACE); RM 102-103',
            'MS River at Bonnet Carre Spillway N (01275, USACE); RM 129-130',
            'MS River at Reserve (01260, USACE); RM 138-139']

staName_short = ['IHNC Surge Barrier East (IHNC-02)',
            'Seabrook Complex (IHNC-01)',
            'Outfall 17th St London Ave Orleans Ave',
            'Lake Pontchartrain at Lakefront Airport',
            'Rigolets near Lake Pontchartrain',
            'Lake Pontchartrain at Mandeville',
            'Bayou Dupre Sector Gate (LPV-144)',
            'Caernarvon Canal Sector Gate',
            'Bayou Bienvenue Floodgate',
            'GIWW at West Closure Complex (WBV-90)',
            'Harvey Canal at Boomtown Casino',
            'Barataria Waterway at Lafitte',
            'Bayou Segnette Closure (WBV-16.2)',
            'Bayou Verret / W. Tie-In Sector Gate (WBV-72/74)',
            'Hero Canal Stop-Log Gate (WBV-09b)',
            'MS River at Empire Floodgate',
            'Baratria Pass at Grand Isle',
            'Houma Navigation Canal',
            'GIWW at Bayou Sale',
            'Bayou Boeuf at Railroad Bridge',
            'Lower Atchafalaya River at Morgan City',
            'Calcasieu River at Cameron',
            'MS River at Venice; RM 10-11',
            'MS River at Empire Lock (NOV-14); RM 29-30',
            'MS River at West Point a la Hache; RM 48-49',
            'MS River at Alliance; RM 62-63',
            'MS River at Carrollton; RM 102-103',
            'MS River at Bonnet Carre Spillway N; RM 129-130',
            'MS River at Reserve; RM 138-139']

# Read run.properties and make a property dictionary
runProp = dict()
#f = open('run.properties','r')
#f = open('AL26_04_nhcConsensus.run.properties','r')
f = open(rprop,'r')
for line in f:
    fields = line.split(':',1)
    try:
        runProp[fields[0].strip()] = fields[1].strip()
    except IndexError:
        continue
f.close()

# Convert advisoryTime to python datetime object
advisory_dt = datetime.strptime(runProp['time.forecast.valid.cdt'],'%Y%m%d%H%M%S')
advisory_dt_long = datetime.strftime(advisory_dt,'%b-%d-%Y %I:%M %p')

#cycleAdvisory_dt = datetime.strptime(runProp['advisory'],'%Y%m%d%H')
cycleAdvisory_dt = datetime.strptime(runProp['forecastValidStart'],'%Y%m%d%H%M%S')
cycleAdvisory_dt_long = datetime.strftime(cycleAdvisory_dt,'%b-%d-%Y %H:%M')

scenario = runProp['asgs.enstorm']
if scenario == 'nhcConsensus':
    scenario_readable = 'NHC Official Track'
else:
    scenario_readable = scenario
if scenario == 'namforecast':
    scenario_readable = 'NAM Forecast'

prs = Presentation('LSU_template.pptx')

slide_layout = prs.slide_layouts[5]
slide_layout_hydro = prs.slide_layouts[6]

numSlides = 1

# Create a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
#if runProp['config.forcing.tropicalcyclone'] != "off": 
if runProp['forcing.tropicalcyclone'] != "off": 
    title.text = runProp['storm class'] + ' ' + runProp['stormname'] + ', ' + scenario_readable + ' Scenario'
    subtitle.text = "Advisory " + runProp['advisory'] + " Issued on " + advisory_dt_long + " CDT" + "\n\n" + \
            'PPT generated on ' + today.strftime('%B %d, %Y %I:%M %p') + ' CDT'
else:
    title.text = scenario_readable + " Cycle Issued on " + cycleAdvisory_dt_long + " UTC"
    subtitle.text = scenario_readable + " Cycle Issued in Local Time on " + advisory_dt_long + " CDT" + "\n\n" + \
            'PPT generated on ' + today.strftime('%B %d, %Y %I:%M %p') + ' CDT'
statement = 'For Official Use Only. Not For Release. \nModel results were produced by the ADCIRC Surge Guidance System (ASGS) and are based on the National Hurricane Center (NHC) forecast track. \nADCIRC-developed hydrographs are an operational planning tool for emergency-response personnel and are not a replacement for National Weather Service (NWS) forecasts.'

fouo = slide.placeholders[10]
fouo.text = statement
numSlides = numSlides + 1
#numSlides = numSlides + 3

# Set slide layout
#left = Inches(1.94)
#top = Inches(1.06)
left = Inches(0.28)
top = Inches(1.13)
#iwidth = Inches(12.32)
#iheight = Inches(4.70)
iwidth = Inches(12.77)
iheight = Inches(5.24)

img_path = fname
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
#if runProp['config.forcing.tropicalcyclone'] != "off": 
if runProp['forcing.tropicalcyclone'] != "off": 
    title.text = 'NHC Advisory ' + runProp['advisory'] + ' ' + scenario_readable + ' Scenario'
else:
    title.text = 'NAM Cycle Forecast Scenario Issued on ' + advisory_dt_long + ' CDT'
#
subtitle.text = "Simulated peak water levels (ft, NAVD88)"
pic = slide.shapes.add_picture(img_path,left,top,width=iwidth,height=iheight)
fouo = slide.placeholders[13]
fouo.text = statement
snum = slide.placeholders[14]
snum.text = str(numSlides)
numSlides = numSlides + 1

############################################################
# TABLE OF CONTENTS SLIDE

slide = prs.slides.add_slide(prs.slide_layouts[7])
title = slide.shapes.title
title.text = 'Table of Contents'
tocText = ''
i = 4
for sta in staName_short:
    tocText = tocText + 'Slide ' + str(i) + ':\t' + sta + '\n'
    i=i+1

maintext = slide.placeholders[1]
maintext.text = tocText

snum = slide.placeholders[14]
snum.text = str(numSlides)
numSlides = numSlides + 1

############################################################

#for shape in slide.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))

#left = Inches(0.42)
#top = Inches(1.15)
left = Inches(0.75)
top = Inches(0.81)
iwidth = Inches(11.84)
iheight = Inches(5.69)

#numSlides = numSlides + 1

i = 0
for image in fnames:
    try: 
        slide = prs.slides.add_slide(slide_layout_hydro)
        title = slide.shapes.title
        title.text = staName[i]
        pic = slide.shapes.add_picture(image,left,top,width=iwidth,height=iheight)
        fouo = slide.placeholders[13]
        fouo.text = statement
        snum = slide.placeholders[14]
        snum.text = str(numSlides)
        numSlides = numSlides + 1
        i = i + 1
    except:
        print("ERROR: buildPPT.py: Could not find " + image + ".")

# Loop through slides
#slides = prs.slides
#for slide in slides:
        #print('slide number %s' % str(slides.index(slide)+1))

#if runProp['config.forcing.tropicalcyclone'] != "off": 
if runProp['forcing.tropicalcyclone'] != "off": 
    pptFile = runProp['stormname'] + "_Adv" + runProp['advisory'] + "_" + scenario + "_" + runProp['forecastValidStart'] + ".pptx"
else:
    #pptFile = runProp['WindModel'] + "_Cycle" + runProp['advisory'] + "_" + scenario + "_" + runProp['forecastValidStart'] + ".pptx"
    pptFile = scenario + "_Cycle_" + runProp['advisory'] + "UTC" + "+gahm.pptx"
prs.save(pptFile)
pFile = open('pptFile.temp','w')
pFile.write(pptFile)
pFile.close()
