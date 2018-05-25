from pptx import Presentation
from pptx.util import Inches


prs = Presentation('LSU_template.pptx')

slide_layout = prs.slide_layouts[5]
slide_layout_hydro = prs.slide_layouts[6]

# Create a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hurricane Isaac"
subtitle.text = "Advisory #25 - Issued on 09:00 UTC Aug-27-2012"

# Set slide layout
left = Inches(1.94)
top = Inches(1.14)

img_path = 'LA_v17a_09_25_nhcOfficial_maxele_0001.jpg'
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "NHC Advisory 25 Consensus"
subtitle.text = "Simulated maximum water level (ft, NAVD88)"
pic = slide.shapes.add_picture(img_path,left,top)

img_path = 'LA_v17a_09_25_veerRight50_maxele_0001.jpg'
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "NHC Advisory 25 veerRight50"
subtitle.text = "Simulated maximum water level (ft, NAVD88)"
pic = slide.shapes.add_picture(img_path,left,top)

left = Inches(0.42)
top = Inches(1.15)

fnames = ['1_WSE_85625.png','2_WSE_76065.png','3_WSE_76030.png','4_WSE_76265.png']
staName = ['17th St. Outfall Canal','Seabrook Complex (IHNC-01)','IHNC Surge Barrier (IHNC-02)',
        'West Closure Complex (WBV-90)']
i = 0
for image in fnames:
    print image
    title.text = staName[i]
    slide = prs.slides.add_slide(slide_layout_hydro)
    title = slide.shapes.title
    pic = slide.shapes.add_picture(image,left,top)
    i = i + 1


# Add image to slide
#img_path = '1_WSE_85625.png'
#slide = prs.slides.add_slide(slide_layout_hydro)
#title = slide.shapes.title
#title.text = "17th St. Outfall Canal"
#pic = slide.shapes.add_picture(img_path,left,top)
#
#img_path = '2_WSE_76065.png'
#slide = prs.slides.add_slide(slide_layout_hydro)
#title = slide.shapes.title
#title.text = "Seabrook Complex (IHNC-01)"
#pic = slide.shapes.add_picture(img_path,left,top)
#
#img_path = '3_WSE_76030.png'
#slide = prs.slides.add_slide(slide_layout_hydro)
#title = slide.shapes.title
#title.text = "IHNC Surge Barrier (IHNC-02)"
#pic = slide.shapes.add_picture(img_path,left,top)
#
#img_path = '4_WSE_76265.png'
#slide = prs.slides.add_slide(slide_layout_hydro)
#title = slide.shapes.title
#title.text = "West Closure Complex (WBV-90)"
#pic = slide.shapes.add_picture(img_path,left,top)


prs.save('Isaac_Adv25_082720180900.pptx')
